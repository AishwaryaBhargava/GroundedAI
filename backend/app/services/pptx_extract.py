from typing import List
import re
import io
from pptx import Presentation


def _normalize_text(s: str) -> str:
    s = s.replace("\u00ad", "")  # soft hyphen
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def extract_pages(pptx_bytes: bytes) -> List[dict]:
    """
    Extract PPTX slides as pages.

    Each slide becomes one page.
    """
    file_like = io.BytesIO(pptx_bytes)
    prs = Presentation(file_like)

    pages: List[dict] = []

    for idx, slide in enumerate(prs.slides):
        texts: List[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)

        page_text = _normalize_text("\n".join(texts))

        pages.append(
            {
                "page": idx + 1,
                "text": page_text,
            }
        )

    return pages
